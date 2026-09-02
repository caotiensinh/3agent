from __future__ import annotations

import csv
import json
import zipfile
from dataclasses import dataclass
from io import BytesIO, StringIO
from pathlib import Path

from defusedxml import ElementTree as SafeElementTree
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps
from pptx import Presentation
from pypdf import PdfReader

try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover - dependency/runtime guard
    fitz = None

MAX_EXTRACTED_CHARS = 200_000
MAX_OOXML_ENTRIES = 512
MAX_OOXML_TOTAL_BYTES = 96 * 1024 * 1024
MAX_OOXML_RATIO = 200
MAX_PDF_PAGES = 300
MAX_PDF_VISION_PAGES = 24
MAX_SHEETS = 32
MAX_ROWS_PER_SHEET = 5_000
MAX_COLS_PER_SHEET = 128
MAX_VISUAL_ASSETS = 32
MAX_VISUAL_PIXELS = 40_000_000
MAX_VISUAL_EDGE = 2048

PLAIN_DOCUMENT_EXTENSIONS = {
    ".csv", ".tsv", ".json", ".jsonl", ".xml", ".yaml", ".yml",
    ".log", ".ini", ".cfg", ".conf",
}
RICH_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
DOCUMENT_EXTENSIONS = PLAIN_DOCUMENT_EXTENSIONS | RICH_DOCUMENT_EXTENSIONS
NATIVE_IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"
}


class DocumentExtractionError(ValueError):
    pass


@dataclass(frozen=True)
class VisualAsset:
    name: str
    locator: str
    media_type: str
    data: bytes
    width: int
    height: int


def _bounded(text: str) -> str:
    return str(text or "").strip()[:MAX_EXTRACTED_CHARS]


def _decode(data: bytes) -> str:
    if b"\x00" in data:
        raise DocumentExtractionError("Text document contains NUL bytes")
    for encoding in ("utf-8-sig", "utf-8", "cp932"):
        try:
            return _bounded(data.decode(encoding))
        except UnicodeDecodeError:
            pass
    return _bounded(data.decode("utf-8", errors="replace"))


def _guard_ooxml(data: bytes) -> None:
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_OOXML_ENTRIES:
                raise DocumentExtractionError("Office document contains too many package entries")
            total = 0
            for info in infos:
                if info.flag_bits & 0x1:
                    raise DocumentExtractionError("Encrypted Office documents are not supported")
                total += max(0, info.file_size)
                if total > MAX_OOXML_TOTAL_BYTES:
                    raise DocumentExtractionError("Office document expands beyond the safety limit")
                if info.compress_size and info.file_size / info.compress_size > MAX_OOXML_RATIO:
                    raise DocumentExtractionError("Office document compression ratio exceeds the safety limit")
    except DocumentExtractionError:
        raise
    except (zipfile.BadZipFile, OSError) as exc:
        raise DocumentExtractionError("Invalid Office Open XML document") from exc


def _extract_pdf(data: bytes) -> tuple[str, list[str]]:
    try:
        reader = PdfReader(BytesIO(data), strict=False)
        if reader.is_encrypted:
            try:
                if reader.decrypt("") == 0:
                    raise DocumentExtractionError("Encrypted PDF is not supported")
            except Exception as exc:
                if isinstance(exc, DocumentExtractionError):
                    raise
                raise DocumentExtractionError("Encrypted PDF is not supported") from exc
        chunks: list[str] = []
        total_chars = 0
        for index, page in enumerate(reader.pages[:MAX_PDF_PAGES], 1):
            text = str(page.extract_text() or "").strip()
            if text:
                block = f"[PDF page {index}]\n{text}"
                chunks.append(block)
                total_chars += len(block)
            if total_chars >= MAX_EXTRACTED_CHARS:
                break
        rendered = _bounded("\n\n".join(chunks))
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(f"PDF parsing failed: {type(exc).__name__}") from exc
    warnings: list[str] = []
    if len(reader.pages) > MAX_PDF_PAGES:
        warnings.append(f"PDF text extraction truncated after {MAX_PDF_PAGES} pages.")
    if not rendered:
        warnings.append("PDF has no machine-readable text; visual analysis is required.")
    return rendered, warnings


def _extract_docx(data: bytes) -> tuple[str, list[str]]:
    _guard_ooxml(data)
    try:
        doc = Document(BytesIO(data))
        rows: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                rows.append(text)
        for table_index, table in enumerate(doc.tables, 1):
            rows.append(f"[Table {table_index}]")
            for row in table.rows:
                cells = [" ".join(cell.text.split()) for cell in row.cells]
                rows.append(" | ".join(cells))
        return _bounded("\n".join(rows)), []
    except Exception as exc:
        raise DocumentExtractionError(f"DOCX parsing failed: {type(exc).__name__}") from exc


def _extract_pptx(data: bytes) -> tuple[str, list[str]]:
    _guard_ooxml(data)
    try:
        deck = Presentation(BytesIO(data))
        rows: list[str] = []
        total_chars = 0
        for index, slide in enumerate(deck.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    texts.append(text)
                if getattr(shape, "has_table", False):
                    table = shape.table
                    for row in table.rows:
                        texts.append(" | ".join(cell.text.strip() for cell in row.cells))
            if texts:
                block = f"[Slide {index}]\n" + "\n".join(texts)
                rows.append(block)
                total_chars += len(block)
            if total_chars >= MAX_EXTRACTED_CHARS:
                break
        return _bounded("\n\n".join(rows)), []
    except Exception as exc:
        raise DocumentExtractionError(f"PPTX parsing failed: {type(exc).__name__}") from exc


def _extract_xlsx(data: bytes) -> tuple[str, list[str]]:
    _guard_ooxml(data)
    try:
        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
        rows: list[str] = []
        warnings: list[str] = []
        total_chars = 0
        try:
            for sheet_index, sheet in enumerate(workbook.worksheets[:MAX_SHEETS], 1):
                heading = f"[Sheet {sheet_index}: {sheet.title}]"
                rows.append(heading)
                total_chars += len(heading)
                for row_index, values in enumerate(
                    sheet.iter_rows(max_col=MAX_COLS_PER_SHEET, values_only=True), 1
                ):
                    if row_index > MAX_ROWS_PER_SHEET:
                        warnings.append(f"Sheet {sheet.title} truncated after {MAX_ROWS_PER_SHEET} rows.")
                        break
                    rendered = ["" if value is None else str(value) for value in values]
                    while rendered and not rendered[-1]:
                        rendered.pop()
                    if rendered:
                        line = "\t".join(rendered)
                        rows.append(line)
                        total_chars += len(line)
                    if total_chars >= MAX_EXTRACTED_CHARS:
                        warnings.append("XLSX text truncated at the global extraction limit.")
                        return _bounded("\n".join(rows)), warnings
            if len(workbook.worksheets) > MAX_SHEETS:
                warnings.append(f"Workbook truncated after {MAX_SHEETS} sheets.")
            return _bounded("\n".join(rows)), warnings
        finally:
            workbook.close()
    except Exception as exc:
        raise DocumentExtractionError(f"XLSX parsing failed: {type(exc).__name__}") from exc


def _extract_csv(data: bytes, delimiter: str | None = None) -> tuple[str, list[str]]:
    text = _decode(data)
    try:
        if delimiter is None:
            try:
                delimiter = csv.Sniffer().sniff(text[:8192], delimiters=",\t;|").delimiter
            except csv.Error:
                delimiter = ","
        reader = csv.reader(StringIO(text), delimiter=delimiter)
        rows: list[str] = []
        warnings: list[str] = []
        total_chars = 0
        for index, row in enumerate(reader, 1):
            if index > MAX_ROWS_PER_SHEET:
                warnings.append(f"Delimited file truncated after {MAX_ROWS_PER_SHEET} rows.")
                break
            line = "\t".join(str(cell) for cell in row[:MAX_COLS_PER_SHEET])
            rows.append(line)
            total_chars += len(line)
            if total_chars >= MAX_EXTRACTED_CHARS:
                warnings.append("Delimited text truncated at the global extraction limit.")
                break
        return _bounded("\n".join(rows)), warnings
    except csv.Error as exc:
        raise DocumentExtractionError("Delimited-text parsing failed") from exc


def _extract_json(data: bytes) -> tuple[str, list[str]]:
    text = _decode(data)
    try:
        value = json.loads(text)
    except json.JSONDecodeError as exc:
        raise DocumentExtractionError("Invalid JSON document") from exc
    return _bounded(json.dumps(value, ensure_ascii=False, indent=2)), []


def _extract_xml(data: bytes) -> tuple[str, list[str]]:
    try:
        root = SafeElementTree.fromstring(data)
    except Exception as exc:
        raise DocumentExtractionError("Invalid or unsafe XML document") from exc
    text = "\n".join(piece.strip() for piece in root.itertext() if piece and piece.strip())
    return _bounded(text), []


def extract_document(filename: str, data: bytes) -> tuple[str, str, list[str]]:
    extension = Path(filename).suffix.casefold()
    if extension == ".pdf":
        text, warnings = _extract_pdf(data)
        kind = "pdf"
    elif extension == ".docx":
        text, warnings = _extract_docx(data)
        kind = "docx"
    elif extension == ".pptx":
        text, warnings = _extract_pptx(data)
        kind = "pptx"
    elif extension == ".xlsx":
        text, warnings = _extract_xlsx(data)
        kind = "xlsx"
    elif extension == ".csv":
        text, warnings = _extract_csv(data)
        kind = "csv"
    elif extension == ".tsv":
        text, warnings = _extract_csv(data, delimiter="\t")
        kind = "tsv"
    elif extension == ".json":
        text, warnings = _extract_json(data)
        kind = "json"
    elif extension == ".xml":
        text, warnings = _extract_xml(data)
        kind = "xml"
    elif extension in {".jsonl", ".yaml", ".yml", ".log", ".ini", ".cfg", ".conf"}:
        text, warnings = _decode(data), []
        kind = extension.lstrip(".")
    else:
        raise DocumentExtractionError(f"Unsupported structured document type: {extension or '<none>'}")
    return _bounded(text), kind, warnings


def _normalize_visual(name: str, locator: str, data: bytes) -> VisualAsset:
    Image.MAX_IMAGE_PIXELS = MAX_VISUAL_PIXELS
    try:
        with Image.open(BytesIO(data)) as source:
            image = ImageOps.exif_transpose(source)
            image.seek(0)
            image.load()
            width, height = image.size
            if width <= 0 or height <= 0 or width * height > MAX_VISUAL_PIXELS:
                raise DocumentExtractionError("Image dimensions exceed the safety limit")
            if max(width, height) > MAX_VISUAL_EDGE:
                image.thumbnail((MAX_VISUAL_EDGE, MAX_VISUAL_EDGE))
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGBA" if "A" in image.mode else "RGB")
            out = BytesIO()
            image.save(out, format="PNG", optimize=True)
            normalized = out.getvalue()
            return VisualAsset(
                name=str(name)[:160],
                locator=str(locator)[:240],
                media_type="image/png",
                data=normalized,
                width=image.width,
                height=image.height,
            )
    except DocumentExtractionError:
        raise
    except Exception as exc:
        raise DocumentExtractionError(f"Unsupported or invalid visual asset: {type(exc).__name__}") from exc


def extract_native_visual(filename: str, data: bytes) -> VisualAsset:
    extension = Path(filename).suffix.casefold()
    if extension not in NATIVE_IMAGE_EXTENSIONS:
        raise DocumentExtractionError(f"Unsupported image type: {extension or '<none>'}")
    return _normalize_visual(filename, "image", data)


def _office_visuals(filename: str, data: bytes, prefix: str) -> tuple[list[VisualAsset], list[str]]:
    _guard_ooxml(data)
    visuals: list[VisualAsset] = []
    warnings: list[str] = []
    with zipfile.ZipFile(BytesIO(data)) as archive:
        for info in archive.infolist():
            if len(visuals) >= MAX_VISUAL_ASSETS:
                warnings.append(f"Visual extraction truncated after {MAX_VISUAL_ASSETS} assets.")
                break
            normalized_name = info.filename.replace("\\", "/")
            if info.is_dir() or not normalized_name.startswith(prefix):
                continue
            try:
                asset = _normalize_visual(
                    Path(normalized_name).name,
                    normalized_name,
                    archive.read(info),
                )
            except DocumentExtractionError:
                warnings.append(f"Unsupported embedded visual skipped: {normalized_name}")
                continue
            visuals.append(asset)
    return visuals, warnings


def _pdf_visuals(filename: str, data: bytes) -> tuple[list[VisualAsset], list[str]]:
    del filename
    if fitz is None:
        return [], ["PyMuPDF is unavailable; PDF visual rendering is disabled."]
    visuals: list[VisualAsset] = []
    warnings: list[str] = []
    try:
        document = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise DocumentExtractionError(f"PDF visual rendering failed: {type(exc).__name__}") from exc
    try:
        page_count = min(document.page_count, MAX_PDF_PAGES)
        candidates: list[int] = []
        for index in range(page_count):
            page = document.load_page(index)
            text = str(page.get_text("text") or "").strip()
            has_images = bool(page.get_images(full=True))
            if len(text) < 160 or has_images:
                candidates.append(index)
        if page_count:
            for index in {0, page_count // 2, page_count - 1}:
                if index not in candidates:
                    candidates.append(index)
        candidates = sorted(candidates)[:MAX_PDF_VISION_PAGES]
        if len(candidates) >= MAX_PDF_VISION_PAGES and page_count > MAX_PDF_VISION_PAGES:
            warnings.append(
                f"PDF visual analysis is bounded to {MAX_PDF_VISION_PAGES} selected pages per upload."
            )
        for index in candidates:
            page = document.load_page(index)
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            visuals.append(
                _normalize_visual(
                    f"page-{index + 1}.png",
                    f"pdf:page:{index + 1}",
                    pix.tobytes("png"),
                )
            )
    finally:
        document.close()
    return visuals, warnings


def extract_visual_assets(filename: str, data: bytes) -> tuple[list[VisualAsset], list[str]]:
    extension = Path(filename).suffix.casefold()
    if extension in NATIVE_IMAGE_EXTENSIONS:
        return [extract_native_visual(filename, data)], []
    if extension == ".pdf":
        return _pdf_visuals(filename, data)
    if extension == ".docx":
        return _office_visuals(filename, data, "word/media/")
    if extension == ".pptx":
        return _office_visuals(filename, data, "ppt/media/")
    if extension == ".xlsx":
        return _office_visuals(filename, data, "xl/media/")
    return [], []
