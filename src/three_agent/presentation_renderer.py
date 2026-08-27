from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PresentationRenderError(RuntimeError):
    pass


class PptxRenderer:
    """Deterministic 16:9 PPTX renderer for validated Agent-2 plans."""

    def __init__(self, font_name: str = "Aptos"):
        self.font_name = font_name

    def _format_run(
        self,
        run,
        size: int,
        *,
        bold: bool = False,
        color: tuple[int, int, int] = (32, 36, 43),
    ) -> None:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = self.font_name
        run.font.color.rgb = RGBColor(*color)

    def _set_title(self, slide, text: str, *, size: int = 30) -> None:
        title = slide.shapes.title
        if title is None:
            raise PresentationRenderError("selected slide layout does not expose a title placeholder")
        title.text = text
        paragraph = title.text_frame.paragraphs[0]
        for run in paragraph.runs:
            self._format_run(run, size, bold=True, color=(18, 46, 86))

    def _add_textbox(
        self,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        size: int,
        *,
        bold: bool = False,
        color: tuple[int, int, int] = (32, 36, 43),
    ):
        shape = slide.shapes.add_textbox(left, top, width, height)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.04)
        frame.margin_right = Inches(0.04)
        frame.margin_top = Inches(0.03)
        frame.margin_bottom = Inches(0.03)
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        self._format_run(run, size, bold=bold, color=color)
        return shape

    def _add_footer(self, slide, slide_no: int, source_ids: list[str]) -> None:
        source_text = "Sources: " + ", ".join(source_ids) if source_ids else ""
        self._add_textbox(
            slide, Inches(0.65), Inches(7.08), Inches(10.9), Inches(0.22),
            source_text, 10, color=(92, 99, 112),
        )
        number = self._add_textbox(
            slide, Inches(12.1), Inches(7.03), Inches(0.6), Inches(0.25),
            str(slide_no), 10, color=(92, 99, 112),
        )
        number.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _add_notes(self, slide, notes: str, source_ids: list[str]) -> None:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is None:
            return
        suffix = ""
        if source_ids:
            suffix = "\nEvidence source IDs: " + ", ".join(source_ids)
        notes_frame.text = (notes or "") + suffix

    @staticmethod
    def _labels(language: str) -> dict[str, str]:
        return {
            "ja": {"fact": "事実", "inference": "推論", "proposal": "提案・制約"},
            "en": {"fact": "Fact", "inference": "Inference", "proposal": "Proposal / limitation"},
            "vi": {"fact": "Sự kiện", "inference": "Suy luận", "proposal": "Đề xuất / giới hạn"},
        }[language]

    def render(self, plan: dict[str, Any], output_path: Path) -> Path:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_layout = prs.slide_layouts[0]
        title_only_layout = prs.slide_layouts[5]
        language = str(plan.get("language", "en"))
        labels = self._labels(language)

        for index, item in enumerate(plan["slides"], start=1):
            is_title = item.get("kind") == "title" and index == 1
            slide = prs.slides.add_slide(title_layout if is_title else title_only_layout)
            title = str(item["title"])
            self._set_title(slide, title, size=36 if is_title else 30)

            if is_title:
                subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                if subtitle is not None:
                    meta_parts = []
                    if plan.get("subtitle"):
                        meta_parts.append(str(plan["subtitle"]))
                    meta_parts.append(f"Audience: {plan['audience']}")
                    meta_parts.append(f"Purpose: {plan['purpose']}")
                    subtitle.text = "\n".join(meta_parts)
                    for paragraph in subtitle.text_frame.paragraphs:
                        for run in paragraph.runs:
                            self._format_run(run, 20, color=(65, 72, 84))
            else:
                y = 1.25
                for claim in item.get("claims", []):
                    prefix = labels["fact"] if claim["kind"] == "verified_fact" else labels["inference"]
                    text = f"{prefix} {claim['claim_id']}: {claim['text']}"
                    self._add_textbox(
                        slide, Inches(0.85), Inches(y), Inches(11.55), Inches(0.78),
                        text, 20,
                    )
                    y += 0.82

                source_slide = item.get("kind") == "sources"
                for text in item.get("context_points", []):
                    height = 0.95 if source_slide else 0.68
                    self._add_textbox(
                        slide, Inches(0.85), Inches(y), Inches(11.55), Inches(height),
                        f"• {text}", 20,
                    )
                    y += height + 0.05

                for text in item.get("proposal_points", []):
                    self._add_textbox(
                        slide, Inches(0.85), Inches(y), Inches(11.55), Inches(0.72),
                        f"{labels['proposal']}: {text}", 20, bold=True, color=(78, 57, 28),
                    )
                    y += 0.77

                if y > 6.82:
                    raise PresentationRenderError(
                        f"slide {index} exceeds safe vertical content budget"
                    )

            source_ids = list(item.get("source_ids", []))
            self._add_footer(slide, index, source_ids)
            self._add_notes(slide, str(item.get("speaker_notes", "")), source_ids)

        prs.save(output_path)
        self.inspect(
            output_path,
            expected_titles=[slide["title"] for slide in plan["slides"]],
        )
        return output_path

    @staticmethod
    def inspect(path: Path, expected_titles: list[str]) -> dict[str, Any]:
        prs = Presentation(path)
        if len(prs.slides) != len(expected_titles):
            raise PresentationRenderError("rendered PPTX slide count mismatch")

        observed: list[str] = []
        title_placeholder_count = 0
        for slide in prs.slides:
            title = slide.shapes.title
            if title is not None:
                title_placeholder_count += 1
                observed.append(title.text.strip())
            else:
                observed.append("")

        if observed != expected_titles:
            raise PresentationRenderError(
                f"rendered PPTX title mismatch: {observed}"
            )
        if title_placeholder_count != len(prs.slides):
            raise PresentationRenderError(
                "every rendered slide must contain a title placeholder"
            )
        return {
            "slide_count": len(prs.slides),
            "titles": observed,
            "title_placeholder_count": title_placeholder_count,
        }


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path | None = None) -> Path:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise PresentationRenderError(
            "PDF output requires LibreOffice/soffice on the host"
        )
    pptx_path = Path(pptx_path).resolve()
    output_dir = Path(output_dir or pptx_path.parent).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        text=True,
        capture_output=True,
        check=False,
        timeout=120,
    )
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if result.returncode != 0 or not pdf_path.exists():
        raise PresentationRenderError(
            f"PDF conversion failed: {result.stderr.strip() or result.stdout.strip()}"
        )
    return pdf_path
