from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XlsxImage
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from three_agent.document_extractors import (
    NATIVE_IMAGE_EXTENSIONS,
    extract_document,
    extract_native_visual,
    extract_visual_assets,
)
from three_agent.knowledge_gateway import KnowledgeGatewayV3
from three_agent.vision import OllamaVisionClient, VisionAnalysis, VisionAnalysisError


class FakeVision:
    model = "fake-vision-v1"

    def __init__(self) -> None:
        self.calls = 0

    def analyze(self, image_bytes: bytes, *, name: str, locator: str) -> VisionAnalysis:
        self.calls += 1
        self.last_bytes = image_bytes
        return VisionAnalysis(
            model=self.model,
            text=(
                "VISIBLE_TEXT\nCamera 12\nIP 192.0.2.12\n"
                "TABLE_OR_VALUES\nPort 554\n"
                f"VISUAL_STRUCTURE\nsource={name} locator={locator}\n"
                "UNCERTAINTIES\nnone"
            ),
        )


def make_png() -> bytes:
    image = Image.new("RGB", (160, 90), "white")
    out = BytesIO()
    image.save(out, format="PNG")
    return out.getvalue()


def make_image(fmt: str) -> bytes:
    image = Image.new("RGB", (96, 64), "white")
    out = BytesIO()
    save_format = "JPEG" if fmt.upper() in {"JPG", "JPEG"} else fmt.upper()
    image.save(out, format=save_format)
    return out.getvalue()


def make_scanned_pdf(image_bytes: bytes) -> bytes:
    out = BytesIO()
    pdf = canvas.Canvas(out, pagesize=(320, 220))
    pdf.drawImage(ImageReader(BytesIO(image_bytes)), 20, 40, width=280, height=158)
    pdf.showPage()
    pdf.save()
    return out.getvalue()


def make_docx(image_bytes: bytes) -> bytes:
    doc = Document()
    doc.add_paragraph("R&D camera status report")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Camera"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "cam12"
    table.cell(1, 1).text = "online"
    doc.add_picture(BytesIO(image_bytes), width=DocxInches(1.2))
    out = BytesIO()
    doc.save(out)
    return out.getvalue()


def make_pptx(image_bytes: bytes) -> bytes:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(5), PptxInches(1))
    box.text_frame.text = "Camera network overview"
    table = slide.shapes.add_table(2, 2, PptxInches(0.5), PptxInches(1.5), PptxInches(4), PptxInches(1.2)).table
    table.cell(0, 0).text = "Node"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "R&D"
    table.cell(1, 1).text = "active"
    slide.shapes.add_picture(BytesIO(image_bytes), PptxInches(5), PptxInches(1.5), width=PptxInches(1.4))
    out = BytesIO()
    deck.save(out)
    return out.getvalue()


def make_xlsx(image_bytes: bytes) -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Status"
    sheet.append(["Camera", "RTSP port"])
    sheet.append(["cam12", 554])
    sheet.add_image(XlsxImage(BytesIO(image_bytes)), "D1")
    out = BytesIO()
    workbook.save(out)
    workbook.close()
    return out.getvalue()


class NativeImageTests(unittest.TestCase):
    def test_common_image_formats_normalize_to_png_for_vision(self) -> None:
        formats = {
            ".png": "PNG",
            ".jpg": "JPEG",
            ".jpeg": "JPEG",
            ".webp": "WEBP",
            ".gif": "GIF",
            ".bmp": "BMP",
            ".tif": "TIFF",
            ".tiff": "TIFF",
        }
        self.assertEqual(set(formats), NATIVE_IMAGE_EXTENSIONS)
        for extension, fmt in formats.items():
            with self.subTest(extension=extension):
                asset = extract_native_visual("sample" + extension, make_image(fmt))
                self.assertEqual(asset.media_type, "image/png")
                self.assertGreater(len(asset.data), 20)
                self.assertEqual((asset.width, asset.height), (96, 64))


class PdfAndOfficeExtractionTests(unittest.TestCase):
    def test_scanned_pdf_has_visual_evidence_even_without_machine_text(self) -> None:
        scanned = make_scanned_pdf(make_png())
        text, kind, warnings = extract_document("scan.pdf", scanned)
        self.assertEqual(kind, "pdf")
        self.assertEqual(text, "")
        self.assertTrue(any("visual analysis" in item for item in warnings))
        visuals, visual_warnings = extract_visual_assets("scan.pdf", scanned)
        self.assertGreaterEqual(len(visuals), 1)
        self.assertEqual(visuals[0].locator, "pdf:page:1")
        self.assertEqual(visuals[0].media_type, "image/png")
        self.assertIsInstance(visual_warnings, list)

    def test_office_files_extract_text_tables_cells_and_embedded_images(self) -> None:
        picture = make_png()
        fixtures = [
            ("report.docx", make_docx(picture), "cam12", "word/media/"),
            ("deck.pptx", make_pptx(picture), "Camera network overview", "ppt/media/"),
            ("status.xlsx", make_xlsx(picture), "554", "xl/media/"),
        ]
        for filename, data, expected_text, expected_locator in fixtures:
            with self.subTest(filename=filename):
                text, kind, _ = extract_document(filename, data)
                self.assertTrue(kind)
                self.assertIn(expected_text, text)
                visuals, _ = extract_visual_assets(filename, data)
                self.assertGreaterEqual(len(visuals), 1)
                self.assertTrue(any(expected_locator in item.locator for item in visuals))


class MultimodalGatewayTests(unittest.TestCase):
    def _gateway(self, root: Path) -> tuple[KnowledgeGatewayV3, FakeVision]:
        gateway = KnowledgeGatewayV3(root, object())
        fake = FakeVision()
        gateway.vision = fake
        return gateway, fake

    def test_native_image_is_understood_and_semantics_are_cached(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            gateway, vision = self._gateway(Path(tmp))
            record = gateway.ingest_upload("camera.jpg", make_image("JPEG"), sender="user-a")
            self.assertEqual(record.kind, "image")
            self.assertEqual(record.image_count, 1)
            context1, diagnostics1 = gateway.build_attachment_context(
                [record.upload_id], "What RTSP port is shown?"
            )
            context2, diagnostics2 = gateway.build_attachment_context(
                [record.upload_id], "Read the camera details again"
            )
            self.assertIn("[LOCAL VISUAL 1:", context1)
            self.assertIn("Port 554", context1)
            self.assertIn("Camera 12", context2)
            self.assertEqual(vision.calls, 1)
            self.assertFalse(any("vision_unavailable" in item for item in diagnostics1 + diagnostics2))

    def test_markdown_html_and_mixed_zip_become_semantic_context(self) -> None:
        picture = make_png()
        office = make_docx(picture)
        archive = BytesIO()
        with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("notes.md", "# Network note\nCamera cam12 uses RTSP.")
            zf.writestr("page.html", "<html><body><h1>R&D</h1><p>Proxy active</p></body></html>")
            zf.writestr("photo.jpg", make_image("JPEG"))
            zf.writestr("office/report.docx", office)
            zf.writestr("nested.zip", b"not-inspected")
        with tempfile.TemporaryDirectory() as tmp:
            gateway, vision = self._gateway(Path(tmp))
            md = gateway.ingest_upload("readme.md", b"# Hello\nlocal document", sender="user-a")
            html = gateway.ingest_upload(
                "page.html", b"<html><body><b>visible html</b><script>ignore()</script></body></html>", sender="user-a"
            )
            mixed = gateway.ingest_upload("bundle.zip", archive.getvalue(), sender="user-a")
            self.assertGreaterEqual(mixed.document_count, 3)
            self.assertGreaterEqual(mixed.image_count, 2)
            context, diagnostics = gateway.build_attachment_context(
                [md.upload_id, html.upload_id, mixed.upload_id],
                "camera proxy RTSP",
            )
            self.assertIn("local document", context)
            self.assertIn("visible html", context)
            self.assertIn("Camera cam12 uses RTSP", context)
            self.assertIn("Proxy active", context)
            self.assertIn("LOCAL VISUAL", context)
            self.assertGreaterEqual(vision.calls, 1)
            self.assertTrue(any("Nested ZIP skipped" in item for item in diagnostics))

    def test_manifest_contains_references_not_vision_response_until_analysis(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            gateway, _ = self._gateway(root)
            record = gateway.ingest_upload("diagram.png", make_png(), sender="user-a")
            manifest_path = root / "uploads" / record.upload_id / "manifest.json"
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            serialized = json.dumps(manifest, ensure_ascii=False)
            self.assertEqual(manifest["schema_version"], 3)
            self.assertIn("visuals", manifest)
            self.assertNotIn("Camera 12", serialized)
            gateway.build_attachment_context([record.upload_id], "read image")
            semantic_files = list((root / "uploads" / record.upload_id / "semantic").glob("vision-*.txt"))
            self.assertEqual(len(semantic_files), 1)
            self.assertIn("Camera 12", semantic_files[0].read_text(encoding="utf-8"))


class VisionBoundaryTests(unittest.TestCase):
    def test_vision_endpoint_is_loopback_only(self) -> None:
        OllamaVisionClient("http://127.0.0.1:11434", "qwen3.6:35b")
        OllamaVisionClient("http://localhost:11434", "qwen3.6:35b")
        with self.assertRaises(VisionAnalysisError):
            OllamaVisionClient("https://example.com", "qwen3.6:35b")


if __name__ == "__main__":
    unittest.main()
