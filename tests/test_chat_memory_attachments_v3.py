from __future__ import annotations

import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from three_agent.chat_attachment_memory import ConversationAttachmentMemory
from three_agent.chat_context import (
    CONTEXT_MODE_CONTINUITY,
    CONTEXT_MODE_FOLLOW_UP,
    build_conversation_context,
)
from three_agent.chat_gateway_v22 import ContinuitySecurityAwareProjectChatService
from three_agent.document_extractors import extract_document
from three_agent.knowledge_gateway import KnowledgeGatewayV2


class ConversationContinuityTests(unittest.TestCase):
    def test_same_conversation_is_injected_without_magic_followup_words(self):
        messages = [
            {"role": "user", "content": "Chúng ta đang thiết kế server AI dùng hai GPU.", "job_id": "a", "status": "completed"},
            {"role": "assistant", "content": "Tôi đề xuất Ubuntu làm server duy nhất trong LAN.", "job_id": "a", "status": "completed"},
            {"role": "user", "content": "Cấu hình firewall nên như thế nào?", "job_id": "b", "status": "completed"},
        ]
        plan = build_conversation_context(messages, "Port 8787 có cần mở toàn mạng không?", current_job_id="current")
        self.assertEqual(plan.mode, CONTEXT_MODE_CONTINUITY)
        self.assertIn("Ubuntu làm server duy nhất", plan.text)
        self.assertGreaterEqual(plan.message_count, 2)

    def test_explicit_reference_keeps_followup_mode(self):
        plan = build_conversation_context(
            [{"role": "assistant", "content": "Phương án A", "job_id": "x", "status": "completed"}],
            "tiếp tục phần trên",
            current_job_id="current",
        )
        self.assertEqual(plan.mode, CONTEXT_MODE_FOLLOW_UP)
        self.assertIn("Phương án A", plan.text)

    def test_current_and_failed_messages_are_not_reinjected(self):
        messages = [
            {"role": "user", "content": "old-safe", "job_id": "old", "status": "completed"},
            {"role": "assistant", "content": "failed-secret", "job_id": "bad", "status": "failed"},
            {"role": "user", "content": "current-duplicate", "job_id": "now", "status": "completed"},
        ]
        plan = build_conversation_context(messages, "new question", current_job_id="now")
        self.assertIn("old-safe", plan.text)
        self.assertNotIn("failed-secret", plan.text)
        self.assertNotIn("current-duplicate", plan.text)


class ConversationAttachmentMemoryTests(unittest.TestCase):
    def test_attachment_references_survive_restart_without_copying_file_content(self):
        with tempfile.TemporaryDirectory() as tmp:
            db = Path(tmp) / "workspace.db"
            memory = ConversationAttachmentMemory(db)
            memory.initialize()
            conversation = "c" * 16
            upload_a = "a" * 16
            upload_b = "b" * 16
            memory.record(conversation, "job-1", [upload_a, upload_b])

            restarted = ConversationAttachmentMemory(db)
            restarted.initialize()
            self.assertEqual(
                restarted.recent_upload_ids(conversation, max_messages=1),
                [upload_a, upload_b],
            )

            with restarted.connect() as conn:
                columns = [
                    str(row[1])
                    for row in conn.execute("PRAGMA table_info(chat_message_attachments)")
                ]
            self.assertEqual(
                columns,
                ["conversation_id", "job_id", "upload_id", "ordinal", "created_at"],
            )

    def test_followup_resolves_recent_attachment_through_owner_validation(self):
        service = object.__new__(ContinuitySecurityAwareProjectChatService)
        service.history = SimpleNamespace(
            get_conversation=lambda owner_key, conversation_id: {
                "conversation_id": conversation_id,
                "messages": [],
            }
        )
        service.attachment_memory = SimpleNamespace(
            recent_upload_ids=lambda conversation_id, **kwargs: ["a" * 16]
        )
        service.orchestrator = SimpleNamespace(knowledge_gateway=object())

        with patch(
            "three_agent.chat_gateway_v22._v4._validate_owned_uploads",
            return_value=["a" * 16],
        ) as validator:
            resolved = service._resolve_submit_uploads(
                "hãy phân tích tiếp file đó",
                channel="web",
                sender="workspace-user:test",
                conversation_id="c" * 16,
                upload_ids=[],
            )
        self.assertEqual(resolved, ["a" * 16])
        validator.assert_called_once()

    def test_unrelated_turn_does_not_silently_inherit_old_attachment(self):
        service = object.__new__(ContinuitySecurityAwareProjectChatService)
        service.history = SimpleNamespace()
        service.attachment_memory = SimpleNamespace()
        service.orchestrator = SimpleNamespace(knowledge_gateway=object())
        resolved = service._resolve_submit_uploads(
            "Hãy giải thích nguyên lý DNS từ đầu.",
            channel="web",
            sender="workspace-user:test",
            conversation_id="c" * 16,
            upload_ids=[],
        )
        self.assertEqual(resolved, [])


class BusinessDocumentExtractionTests(unittest.TestCase):
    def test_csv_docx_pptx_xlsx_pdf_are_readable(self):
        text, kind, _ = extract_document("sample.csv", b"name,value\nalpha,42\n")
        self.assertEqual(kind, "csv")
        self.assertIn("alpha", text)

        buf = BytesIO()
        doc = Document()
        doc.add_paragraph("DOCX evidence 123")
        doc.save(buf)
        text, kind, _ = extract_document("sample.docx", buf.getvalue())
        self.assertEqual(kind, "docx")
        self.assertIn("DOCX evidence 123", text)

        buf = BytesIO()
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        box.text_frame.text = "PPTX evidence 456"
        deck.save(buf)
        text, kind, _ = extract_document("sample.pptx", buf.getvalue())
        self.assertEqual(kind, "pptx")
        self.assertIn("PPTX evidence 456", text)

        buf = BytesIO()
        book = Workbook()
        sheet = book.active
        sheet["A1"] = "XLSX evidence"
        sheet["B1"] = 789
        book.save(buf)
        text, kind, _ = extract_document("sample.xlsx", buf.getvalue())
        self.assertEqual(kind, "xlsx")
        self.assertIn("XLSX evidence", text)

        buf = BytesIO()
        pdf = canvas.Canvas(buf)
        pdf.drawString(72, 720, "PDF evidence 999")
        pdf.save()
        text, kind, _ = extract_document("sample.pdf", buf.getvalue())
        self.assertEqual(kind, "pdf")
        self.assertIn("PDF evidence 999", text)

    def test_gateway_turns_business_document_into_local_source(self):
        with tempfile.TemporaryDirectory() as tmp:
            gateway = KnowledgeGatewayV2(Path(tmp), object())
            record = gateway.ingest_upload(
                "notes.csv",
                b"topic,result\nnetwork,pass\n",
                sender="workspace-user:test",
            )
            sources, diagnostics = gateway.load_upload_sources([record.upload_id])
            self.assertFalse(diagnostics)
            self.assertEqual(len(sources), 1)
            self.assertIn("network", sources[0].extracted_text)
            self.assertTrue(sources[0].url.startswith("upload://"))

    def test_long_attachment_retrieval_finds_relevant_text_beyond_old_12k_prefix(self):
        with tempfile.TemporaryDirectory() as tmp:
            gateway = KnowledgeGatewayV2(Path(tmp), object())
            filler = ("ordinary filler line without the target term\n" * 600).encode("utf-8")
            marker = "DEVICE-ALPHA-9876 requires LAN port 8787 only\n".encode("utf-8")
            record = gateway.ingest_upload(
                "long-notes.log",
                filler + marker,
                sender="workspace-user:test",
            )
            context, diagnostics = gateway.build_attachment_context(
                [record.upload_id],
                "What does DEVICE-ALPHA-9876 require?",
                max_chars=24_000,
            )
            self.assertNotIn("image_not_semantically_parsed", " ".join(diagnostics))
            self.assertIn("DEVICE-ALPHA-9876", context)
            self.assertIn("8787", context)
            self.assertLessEqual(len(context), 24_000)


if __name__ == "__main__":
    unittest.main()
