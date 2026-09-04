from __future__ import annotations

import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from three_agent.document_extractors import extract_document
from three_agent.knowledge_gateway import EXTENDED_UPLOAD_EXTENSIONS, KnowledgeGatewayV2
from three_agent.workspace_frontend import WORKSPACE_HTML


class BusinessDocumentExtractionTests(unittest.TestCase):
    def test_csv_docx_pptx_xlsx_pdf_are_readable(self):
        text, kind, _ = extract_document("sample.csv", b"name,value\nalpha,42\n")
        self.assertEqual(kind, "csv")
        self.assertIn("alpha", text)

        buf = BytesIO()
        doc = Document()
        doc.add_paragraph("DOCX evidence 123")
        doc.save(buf)
        text, kind, _ = extract_document("sample.docx", buf.getvalue())
        self.assertEqual(kind, "docx")
        self.assertIn("DOCX evidence 123", text)

        buf = BytesIO()
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        box.text_frame.text = "PPTX evidence 456"
        deck.save(buf)
        text, kind, _ = extract_document("sample.pptx", buf.getvalue())
        self.assertEqual(kind, "pptx")
        self.assertIn("PPTX evidence 456", text)

        buf = BytesIO()
        book = Workbook()
        sheet = book.active
        sheet["A1"] = "XLSX evidence"
        sheet["B1"] = 789
        book.save(buf)
        text, kind, _ = extract_document("sample.xlsx", buf.getvalue())
        self.assertEqual(kind, "xlsx")
        self.assertIn("XLSX evidence", text)

        buf = BytesIO()
        pdf = canvas.Canvas(buf)
        pdf.drawString(72, 720, "PDF evidence 999")
        pdf.save()
        text, kind, _ = extract_document("sample.pdf", buf.getvalue())
        self.assertEqual(kind, "pdf")
        self.assertIn("PDF evidence 999", text)

    def test_long_attachment_retrieval_finds_relevant_text_beyond_old_12k_prefix(self):
        with tempfile.TemporaryDirectory() as tmp:
            gateway = KnowledgeGatewayV2(Path(tmp), object())
            filler = (
                "ordinary filler line without the target term\n" * 600
            ).encode("utf-8")
            marker = "DEVICE-ALPHA-9876 requires LAN port 8787 only\n".encode("utf-8")
            record = gateway.ingest_upload(
                "long-notes.log",
                filler + marker,
                sender="workspace-user:test",
            )
            context, diagnostics = gateway.build_attachment_context(
                [record.upload_id],
                "What does DEVICE-ALPHA-9876 require?",
                max_chars=24_000,
            )
            self.assertNotIn(
                "image_not_semantically_parsed",
                " ".join(diagnostics),
            )
            self.assertIn("DEVICE-ALPHA-9876", context)
            self.assertIn("8787", context)
            self.assertLessEqual(len(context), 24_000)

    def test_business_document_capability_and_local_processing_ui_are_exposed(self):
        self.assertTrue(
            {".pdf", ".docx", ".pptx", ".xlsx", ".csv"}.issubset(
                EXTENDED_UPLOAD_EXTENSIONS
            )
        )
        self.assertIn("Uploading and processing files locally", WORKSPACE_HTML)
        self.assertIn("Attachments processed locally", WORKSPACE_HTML)


if __name__ == "__main__":
    unittest.main()
