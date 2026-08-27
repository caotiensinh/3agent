import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from three_agent.artifacts import ArtifactManager
from three_agent.presentation_model import (
    EvidenceCatalog,
    PresentationOptions,
    PresentationValidationError,
    build_dry_run_plan,
    handoff_is_presentable,
    normalize_plan,
)
from three_agent.presentation_renderer import PptxRenderer


class PresentationAgentContractTests(unittest.TestCase):
    def handoff_payload(self):
        return {
            "schema_version": "1.0",
            "task_id": "TASK-1",
            "presentation_ready": True,
            "blockers": [],
            "key_facts": [
                {"fact_id": "F001", "claim": "Product A supports feature X.", "source_ids": ["S1"], "confidence": "medium"},
                {"fact_id": "F002", "claim": "Product B costs more than Product A.", "source_ids": ["S2"], "confidence": "medium"},
            ],
            "inferences": [
                {"claim": "Product A may fit the target use case better.", "source_ids": ["S1", "S2"], "confidence": "high"}
            ],
            "conflicts": [],
            "unresolved_items": ["Long-term support terms are not verified."],
            "conclusion": "Evidence supports a controlled comparison.",
            "recommended_next_actions": ["Run a two-week proof of concept."],
            "sources": [
                {"source_id": "S1", "title": "Vendor A", "url": "https://a.example", "fetch_status": "ok"},
                {"source_id": "S2", "title": "Vendor B", "url": "https://b.example", "fetch_status": "ok"},
            ],
            "quality_metrics": {"usable_source_count": 2, "verified_fact_count": 2},
        }

    def valid_raw_plan(self):
        return {
            "title": "Decision brief",
            "subtitle": "R&D review",
            "slides": [
                {"kind": "title", "title": "Decision brief", "claim_refs": [], "proposal_points": [], "context_points": ["Internal evaluation"], "speaker_notes": "Open with the decision objective."},
                {"kind": "comparison", "title": "Verified comparison", "claim_refs": ["F001", "F002"], "proposal_points": [], "context_points": [], "speaker_notes": "Explain only the two verified points."},
                {"kind": "decision", "title": "Recommendation", "claim_refs": ["I001"], "proposal_points": ["Run a two-week proof of concept."], "context_points": [], "speaker_notes": "Label the inference and proposal distinctly."},
            ],
        }

    def test_handoff_catalog_preserves_fact_ids_truth_classes_and_sources(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        claims = catalog.claim_map
        self.assertEqual(claims["F001"].kind, "verified_fact")
        self.assertEqual(claims["I001"].kind, "inference")
        self.assertEqual(claims["F001"].source_ids, ("S1",))
        self.assertEqual(claims["I001"].source_ids, ("S1", "S2"))

    def test_handoff_gate_requires_ready_and_verified_fact(self):
        ready, _ = handoff_is_presentable(self.handoff_payload())
        self.assertTrue(ready)
        blocked = self.handoff_payload()
        blocked["presentation_ready"] = False
        blocked["blockers"] = ["CRITICAL_SOURCE_CONFLICT"]
        ready, reason = handoff_is_presentable(blocked)
        self.assertFalse(ready)
        self.assertIn("CRITICAL_SOURCE_CONFLICT", reason)

    def test_unknown_claim_reference_is_hard_failure(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        plan = self.valid_raw_plan()
        plan["slides"][1]["claim_refs"] = ["F999"]
        with self.assertRaises(PresentationValidationError):
            normalize_plan(plan, catalog, PresentationOptions(), "Fallback")

    def test_inference_only_deck_is_hard_failure(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        plan = self.valid_raw_plan()
        plan["slides"][1]["claim_refs"] = ["I001"]
        plan["slides"][2]["claim_refs"] = ["I001"]
        with self.assertRaisesRegex(PresentationValidationError, "at least one verified fact"):
            normalize_plan(plan, catalog, PresentationOptions(), "Fallback")

    def test_duplicate_titles_are_hard_failure(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        plan = self.valid_raw_plan()
        plan["slides"][2]["title"] = plan["slides"][1]["title"]
        with self.assertRaises(PresentationValidationError):
            normalize_plan(plan, catalog, PresentationOptions(), "Fallback")

    def test_materialized_claims_are_exact_and_appendices_are_deterministic(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        plan, qa = normalize_plan(self.valid_raw_plan(), catalog, PresentationOptions(slide_count=6), "Fallback")
        self.assertEqual(qa["status"], "pass")
        self.assertTrue(qa["visible_facts_source_bounded"])
        self.assertTrue(any(slide["kind"] == "sources" for slide in plan["slides"]))
        self.assertTrue(any(slide["kind"] == "limitations" for slide in plan["slides"]))
        factual = next(s for s in plan["slides"] if s["title"] == "Verified comparison")
        self.assertEqual(factual["claims"][0]["text"], "Product A supports feature X.")
        self.assertEqual(factual["source_ids"], ["S1", "S2"])

    def test_dry_run_is_deterministic_and_source_bounded(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        options = PresentationOptions(language="ja", slide_count=6)
        a = build_dry_run_plan("Title", "Request", options, catalog)
        b = build_dry_run_plan("Title", "Request", options, catalog)
        self.assertEqual(a, b)
        self.assertEqual(next(s for s in a["slides"] if s["kind"] == "content")["claims"][0]["text"], "Product A supports feature X.")

    def test_pptx_reopens_with_title_placeholders_and_notes(self):
        catalog = EvidenceCatalog.from_handoff(self.handoff_payload())
        plan, _ = normalize_plan(self.valid_raw_plan(), catalog, PresentationOptions(slide_count=6), "Fallback")
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "deck.pptx"
            PptxRenderer().render(plan, path)
            self.assertTrue(path.exists())
            prs = Presentation(path)
            self.assertEqual(len(prs.slides), len(plan["slides"]))
            self.assertTrue(all(slide.shapes.title is not None for slide in prs.slides))
            first_notes = prs.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("decision objective", first_notes)
            factual_slide = next(slide for slide in prs.slides if slide.shapes.title.text == "Verified comparison")
            self.assertIn("Evidence source IDs: S1, S2", factual_slide.notes_slide.notes_text_frame.text)

    def test_artifact_manager_finds_latest_cross_day_handoff(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            older = root / "research" / "2026-08-26"
            newer = root / "research" / "2026-08-27"
            older.mkdir(parents=True)
            newer.mkdir(parents=True)
            old_path = older / "TASK-1_handoff.json"
            new_path = newer / "TASK-1_handoff.json"
            old_path.write_text(json.dumps({"v": 1}), encoding="utf-8")
            new_path.write_text(json.dumps({"v": 2}), encoding="utf-8")
            manager = ArtifactManager(root)
            found = manager.find_latest_task_artifact("research", "TASK-1", suffix="_handoff.json")
            self.assertEqual(found, new_path)


if __name__ == "__main__":
    unittest.main()
